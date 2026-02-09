from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_chaincare_presentation():
    """Create ChainCare PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    PRIMARY_BLUE = RGBColor(41, 98, 255)
    TEAL = RGBColor(0, 150, 136)
    DARK_GRAY = RGBColor(33, 33, 33)
    LIGHT_GRAY = RGBColor(242, 242, 242)
    WHITE = RGBColor(255, 255, 255)
    RED = RGBColor(244, 67, 54)
    GREEN = RGBColor(76, 175, 80)
    
    def add_title_slide(title, subtitle, team_members=None):
        """Create title slide with gradient background"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add background gradient (simulated with rectangle)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = PRIMARY_BLUE
        bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.2), Inches(8), Inches(0.6)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = WHITE
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Team members
        if team_members:
            team_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(4.5), Inches(7), Inches(1.5)
            )
            team_frame = team_box.text_frame
            for i, member in enumerate(team_members):
                if i > 0:
                    team_frame.add_paragraph()
                p = team_frame.paragraphs[i]
                p.text = member
                p.font.size = Pt(16)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
        
        # Footer
        footer_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(8), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_para = footer_frame.paragraphs[0]
        footer_para.text = "Final Year Project | Computer Science & Engineering | January 2026"
        footer_para.font.size = Pt(14)
        footer_para.font.color.rgb = WHITE
        footer_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, content_items, notes=""):
        """Create content slide with title and bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.fill.background()
        
        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, prs.slide_width, Inches(1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = PRIMARY_BLUE
        title_bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            
            # Check for bullet level (based on indentation in text)
            if item.startswith("    "):
                p.text = item.strip()
                p.level = 1
                p.font.size = Pt(16)
            elif item.startswith("  "):
                p.text = item.strip()
                p.level = 1
                p.font.size = Pt(18)
            else:
                p.text = item
                p.level = 0
                p.font.size = Pt(20)
            
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(6)
        
        # Add notes
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
        
        return slide
    
    def add_two_column_slide(title, left_content, right_content, notes=""):
        """Create slide with two columns"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.fill.background()
        
        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, prs.slide_width, Inches(1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = TEAL
        title_bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        for i, item in enumerate(left_content):
            if i > 0:
                left_frame.add_paragraph()
            p = left_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
        
        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        for i, item in enumerate(right_content):
            if i > 0:
                right_frame.add_paragraph()
            p = right_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
        
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
        
        return slide
    
    # SLIDE 1: Title Slide
    add_title_slide(
        "ChainCare",
        "Blockchain-Secured Medical Records Management System",
        [
            "Team Members:",
            "• [Member 1] - Authentication & User Management",
            "• [Member 2] - Medical Records & Data Management",
            "• [Member 3] - Blockchain Audit Trail & Verification",
            "• [Member 4] - AI Clinical Assistants & Doctor Workflow"
        ]
    )
    
    # SLIDE 2: Problem Statement
    add_content_slide(
        "Problem Statement",
        [
            "❌ DATA FRAGMENTATION",
            "  Patient reports scattered across hospitals",
            "  No centralized access for emergencies",
            "",
            "❌ SECURITY CONCERNS",
            "  Paper records easily lost/tampered",
            "  Digital records vulnerable to unauthorized access",
            "",
            "❌ TRUST ISSUES",
            "  Cannot verify if records were altered",
            "  No audit trail of who accessed data when",
            "",
            "❌ INFORMATION ASYMMETRY",
            "  Patients don't understand medical jargon",
            "  Doctors spend time re-diagnosing"
        ],
        notes="Real-world examples: lost prescriptions during emergencies, data breaches in hospitals"
    )
    
    # SLIDE 3: Our Solution
    add_content_slide(
        "Our Solution - ChainCare",
        [
            "A Complete Medical Records Ecosystem",
            "",
            "✅ SECURE STORAGE",
            "  Cloud-based (Firebase) accessible anywhere",
            "  AES-256 encryption at rest and in transit",
            "",
            "✅ BLOCKCHAIN AUDIT TRAIL",
            "  Immutable record of all data access",
            "  Polygon blockchain verification",
            "",
            "✅ AI-POWERED ASSISTANCE",
            "  Doctor Clinical Assistant (data analysis)",
            "  Patient Medical Assistant (education)",
            "",
            "✅ SMART CONSENT",
            "  QR-based patient consent",
            "  Every access logged and traceable"
        ],
        notes="Emphasize multi-faceted approach - not just storage, but complete ecosystem"
    )
    
    # SLIDE 4: System Architecture
    add_content_slide(
        "System Architecture",
        [
            "4-Layer Architecture",
            "",
            "Layer 1: FRONTEND",
            "  Flutter Mobile App (iOS & Android)",
            "",
            "Layer 2: BACKEND - Firebase Ecosystem",
            "  • Authentication (Email, OAuth, Phone OTP)",
            "  • Firestore Database (NoSQL documents)",
            "  • Cloud Storage (medical documents)",
            "  • Cloud Functions (serverless Node.js)",
            "",
            "Layer 3: BLOCKCHAIN - Polygon Amoy Testnet",
            "  • Smart Contract (Solidity)",
            "  • Daily audit hash posting",
            "",
            "Layer 4: AI - Google Gemini 2.0 Flash",
            "  • RAG architecture for context"
        ],
        notes="Explain bidirectional data flow, emphasize serverless nature"
    )
    
    # SLIDE 5: Technology Stack
    add_content_slide(
        "Technology Stack",
        [
            "Frontend: Flutter 3.24",
            "  Cross-platform (1 codebase for iOS + Android)",
            "",
            "Authentication: Firebase Auth",
            "  Built-in email, OAuth, phone OTP",
            "",
            "Database: Cloud Firestore",
            "  Real-time sync, offline support",
            "",
            "Blockchain: Polygon Amoy",
            "  Low gas fees ($0.004/tx vs Ethereum $50+)",
            "",
            "Smart Contract: Solidity 0.8.0",
            "  Ethereum-compatible, battle-tested",
            "",
            "AI: Google Gemini 2.0 Flash",
            "  Free tier, medical knowledge, 32K context",
            "",
            "Backend: Firebase Cloud Functions",
            "  Serverless, auto-scaling, Node.js"
        ],
        notes="Cost efficiency: runs on Firebase free tier for testing, <$50/month for 10K users"
    )
    
    # SLIDE 6: User Roles & Journeys
    add_two_column_slide(
        "User Roles & Journeys",
        [
            "🔵 PATIENT",
            "• Sign up with email/Google",
            "• Complete health profile",
            "• Upload medical records",
            "• Generate QR code",
            "• Chat with AI Assistant",
            "• View audit trail",
            "",
            "🟢 DOCTOR",
            "• Sign up and submit credentials",
            "• Wait for admin verification",
            "• Scan patient QR for consent",
            "• View patient history",
            "• Chat with Clinical Assistant",
            "• Upload prescriptions/notes"
        ],
        [
            "🟣 ADMIN",
            "• Review doctor applications",
            "• Approve/reject after ML analysis",
            "• Monitor blockchain posting",
            "• View global audit logs",
            "• System health dashboard"
        ],
        notes="Quick walkthrough of typical user flows"
    )
    
    # SLIDE 7: Biometric App Lock
    add_content_slide(
        "Feature: Biometric App Lock",
        [
            "Banking-Grade App Security",
            "",
            "TRIGGER CONDITIONS:",
            "  ✅ Cold Start (first launch after device restart)",
            "  ✅ Resume from Background (>60 seconds minimized)",
            "  ❌ Does NOT trigger on internal navigation",
            "",
            "PLATFORM SUPPORT:",
            "  • iOS: FaceID, TouchID",
            "  • Android: Fingerprint, Face Unlock",
            "  • Fallback: Device passcode",
            "",
            "SECURITY:",
            "  Biometric data NEVER leaves device (OS secure enclave)",
            "",
            "IMPLEMENTATION:",
            "  • Lifecycle management with AppLifecycleState",
            "  • Time-based background detection",
            "  • Platform-specific biometric APIs"
        ],
        notes="Demo live biometric lock if possible on presentation machine"
    )
    
    # SLIDE 8: QR-Based Consent
    add_content_slide(
        "Feature: QR-Based Consent",
        [
            "Explicit Patient Consent Mechanism",
            "",
            "PROCESS FLOW:",
            "  Step 1: Patient generates Medical ID with QR code",
            "  Step 2: Doctor scans QR with in-app scanner",
            "  Step 3: System creates audit log entry",
            "  Step 4: Doctor gains READ access to patient records",
            "  Step 5: Audit entry included in blockchain post",
            "",
            "QR DATA:",
            "  Only contains patient UID (28-char Firebase ID)",
            "",
            "SECURITY:",
            "  Requires doctor authentication + QR physical proximity",
            "",
            "AUDIT LOGGING:",
            "  Every scan creates immutable audit trail",
            "  (action, timestamp, doctor ID, patient ID)"
        ],
        notes="Contrast with password sharing (insecure, no audit)"
    )
    
    # SLIDE 9: ML Kit Document Verification
    add_content_slide(
        "Feature: AI-Assisted Doctor Verification",
        [
            "ML-Powered Credential Verification",
            "",
            "PROCESS:",
            "  1. Doctor uploads medical license image",
            "  2. ML Kit OCR extracts text on-device",
            "  3. AI analyzes:",
            "     • Name matching (vs profile name)",
            "     • Medical keywords detection",
            "     • License number patterns",
            "     • Document quality assessment",
            "  4. Confidence score generated (0.0 - 1.0)",
            "  5. Admin reviews with AI recommendations",
            "",
            "CONFIDENCE THRESHOLDS:",
            "  ✅ 0.9-1.0: High (green) - Fast-track approval",
            "  ✅ 0.7-0.9: Good (green) - Likely approved",
            "  ⚠️  0.4-0.7: Low (yellow) - Careful review",
            "  ❌ 0.0-0.4: Reject (red) - Cannot submit"
        ],
        notes="Emphasize: AI ASSISTS admin, doesn't replace human judgment"
    )
    
    # SLIDE 10: Blockchain Audit Trail
    add_content_slide(
        "Blockchain Audit Trail - Architecture",
        [
            "Two-Layer Tamper-Evident System",
            "",
            "LAYER 1: Real-Time Hash Chain (Firestore)",
            "  • Every sensitive action creates audit entry",
            "  • Each entry links to previous via SHA-256 hash",
            "  • Breaking chain immediately detectable",
            "",
            "LAYER 2: Daily Blockchain Anchoring (Polygon)",
            "  • Midnight UTC: aggregate all day's entries",
            "  • Compute Merkle root (1000 entries → 1 hash)",
            "  • Post to public Polygon blockchain",
            "  • External, immutable proof",
            "",
            "WHY BOTH LAYERS?",
            "  • Real-time: Immediate logging (no delay)",
            "  • Blockchain: Prevents retroactive tampering",
            "    (even by Firebase admin)"
        ],
        notes="Analogy: hash chain = sealed envelope, blockchain = public notarization"
    )
    
    # SLIDE 11: Merkle Tree Explained
    add_content_slide(
        "Blockchain: Merkle Tree Explained",
        [
            "Efficient Aggregation with Merkle Trees",
            "",
            "BOTTOM LEVEL (Leaves):",
            "  Hash(Audit #1) | Hash(Audit #2) | Hash(Audit #3) | Hash(Audit #4)",
            "",
            "LEVEL 2:",
            "  Hash(H1 + H2) | Hash(H3 + H4)",
            "",
            "LEVEL 3 (Root):",
            "  Hash(H(1+2) + H(3+4)) ← Posted to blockchain",
            "",
            "TAMPERING DETECTION:",
            "  • Change Audit #2 → H2 changes",
            "  • H2 changes → H(1+2) changes",
            "  • H(1+2) changes → Root changes",
            "  • Root on blockchain doesn't match → TAMPER DETECTED!",
            "",
            "EFFICIENCY:",
            "  Logarithmic scaling (1M entries → 20 levels)"
        ],
        notes="Emphasize logarithmic efficiency for large datasets"
    )
    
    # SLIDE 12: Smart Contract
    add_content_slide(
        "Blockchain: Smart Contract",
        [
            "Solidity Smart Contract on Polygon",
            "",
            "CONTRACT STRUCTURE:",
            "  • AuditEntry struct (date, merkleRoot, count, timestamp)",
            "  • Array of all audit entries",
            "  • storeAuditHash() function",
            "  • getAudit() view function",
            "",
            "DEPLOYMENT:",
            "  • Contract Address: 0x56bBF330d155B30aAeb904B93D21EeBCb1f96aB6",
            "  • Network: Polygon Amoy Testnet (Chain ID: 80002)",
            "  • Explorer: amoy.polygonscan.com",
            "",
            "VALIDATION:",
            "  • Requires 64-character SHA-256 hash",
            "  • Immutable once posted",
            "  • Publicly verifiable",
            "",
            "COST:",
            "  ~$0.004 USD per transaction"
        ],
        notes="Show actual transaction on PolygonScan if internet available"
    )
    
    # SLIDE 13: AI RAG Architecture
    add_content_slide(
        "AI Clinical Assistants - RAG Architecture",
        [
            "RAG: Retrieval Augmented Generation",
            "",
            "THE PROBLEM with Traditional LLMs:",
            "  • No knowledge of patient-specific records",
            "  • Hallucinates data",
            "  • Cannot answer 'What medications am I on?'",
            "",
            "RAG SOLUTION:",
            "  1. RETRIEVE patient records from Firestore",
            "  2. FORMAT into readable context (categorized)",
            "  3. AUGMENT AI prompt with records",
            "  4. GENERATE response grounded in actual data",
            "",
            "CONTEXT EXAMPLE:",
            "  LAB TESTS (15 records)",
            "    Record #1 - Jan 09, 2025",
            "    Type: Blood Test - CBC",
            "    Hemoglobin: 13.5 g/dL, WBC: 7000/μL",
            "",
            "BENEFITS:",
            "  • No hallucination (cites actual records)",
            "  • Up-to-date (fetches latest data)",
            "  • Context-aware (analyzes patterns)"
        ],
        notes="Example: 'Is patient anemic?' → AI cites Record #1 with actual values"
    )
    
    # SLIDE 14: AI Dual Personalities
    add_two_column_slide(
        "AI - Dual Personalities",
        [
            "DOCTOR CLINICAL ASSISTANT",
            "",
            "Tone: Professional, concise",
            "",
            "Language: Medical terminology",
            "",
            "Citations: 'According to Record #3, Dec 15...'",
            "",
            "Focus: Data analysis, trends",
            "",
            "Example:",
            "'Mild anemia detected (Hgb 12.5, Record #2). Appears chronic based on previous CBC.'"
        ],
        [
            "PATIENT MEDICAL ASSISTANT",
            "",
            "Tone: Warm, conversational",
            "",
            "Language: Simple analogies",
            "",
            "Citations: 'Your Dec 15 blood test showed...'",
            "",
            "Focus: Education, reassurance",
            "",
            "Example:",
            "'Your Dec test shows mild anemia - your red blood cell count is a bit low. Think of it like fewer delivery trucks carrying oxygen. Common and treatable!'"
        ],
        notes="Demo both if time permits (show different responses to same question)"
    )
    
    # SLIDE 15: Security Pyramid
    add_content_slide(
        "Security Pyramid - Defense in Depth",
        [
            "6 Layers of Security",
            "",
            "Layer 1 (Foundation): DEVICE BIOMETRIC",
            "  FaceID, TouchID, Fingerprint - Data never leaves device",
            "",
            "Layer 2: FIREBASE AUTHENTICATION",
            "  Email/password (bcrypt), Google OAuth, Phone OTP",
            "",
            "Layer 3: DATA ENCRYPTION",
            "  At-rest: AES-256 | In-transit: HTTPS/TLS 1.3",
            "",
            "Layer 4: ROLE-BASED ACCESS CONTROL",
            "  Firestore Security Rules (server-enforced)",
            "",
            "Layer 5: AUDIT LOGGING",
            "  Every action logged, hash chain prevents tampering",
            "",
            "Layer 6 (Apex): BLOCKCHAIN VERIFICATION",
            "  Daily Merkle root on Polygon - Public, immutable"
        ],
        notes="Defense in depth: even if one layer breached, others protect data"
    )
    
    # SLIDE 16: Firestore Security Rules
    add_content_slide(
        "Firestore Security Rules",
        [
            "Server-Side Enforcement (Cannot Be Bypassed)",
            "",
            "MEDICAL RECORDS:",
            "  ✅ Patient can CRUD own records",
            "  ✅ Doctors can READ any patient's records",
            "  ✅ Doctors can WRITE records (prescriptions)",
            "",
            "AUDIT CHAIN:",
            "  ✅ Users can READ audit logs",
            "  ❌ Only Cloud Functions can WRITE",
            "",
            "CHAT HISTORY:",
            "  ✅ Patient can access own chat history",
            "  ❌ No one else (including doctors/admin)",
            "",
            "ATTACK PREVENTION:",
            "  • Patient tries to read another's record → DENIED",
            "  • User tries to change own role → DENIED",
            "  • Hacker modifies APK to write audit → DENIED",
            "",
            "Runs on Google's servers, not in app!"
        ],
        notes="Stress this runs on Google's servers, cannot be bypassed by modified APK"
    )
    
    # SLIDE 17: Performance & Optimization
    add_content_slide(
        "Performance & Optimization",
        [
            "Optimization Strategies Implemented",
            "",
            "1. FIRESTORE QUERY OPTIMIZATION",
            "   • Indexing on timestamp field",
            "   • Pagination (limit 50 records per fetch)",
            "   • Offline persistence (instant load from cache)",
            "",
            "2. IMAGE COMPRESSION",
            "   Before: 5 MB JPEG → After: 800 KB JPEG (84% reduction)",
            "",
            "3. AI CONTEXT PRUNING",
            "   • Last 20 messages in conversation history",
            "   • Truncate OCR text to 500 chars per record",
            "   • Select 50 most relevant from >100 records",
            "",
            "4. APP SIZE REDUCTION",
            "   45 MB → 32 MB APK (28% reduction)",
            "",
            "BENCHMARKS:",
            "  • App launch: 1.2s (cold start)",
            "  • Record upload: 3-5s (5 MB PDF)",
            "  • AI response: 1-2s (average)",
            "  • QR scan: <1s"
        ],
        notes="Emphasize real-world usability over theoretical performance"
    )
    
    # SLIDE 18: Cost Analysis
    add_content_slide(
        "Cost Analysis",
        [
            "Development vs Production Costs",
            "",
            "DEVELOPMENT (Testnet/Free):",
            "  All services: $0/month",
            "",
            "PRODUCTION (10,000 users):",
            "  • Firebase Auth: $0 (unlimited)",
            "  • Firestore: ~$30/month",
            "  • Cloud Storage: ~$10/month",
            "  • Cloud Functions: ~$5/month",
            "  • Polygon Blockchain: ~$1.50/month",
            "  • Gemini AI: ~$1/month",
            "  TOTAL: ~$47.50/month",
            "",
            "SCALING:",
            "  • 100 users: $0 (within free tier)",
            "  • 1,000 users: ~$15/month",
            "  • 10,000 users: ~$47/month",
            "  • 100,000 users: ~$250/month",
            "",
            "Per-user cost at 10K scale: $0.00475/month"
        ],
        notes="Highlight serverless cost advantages (pay only for usage, not idle capacity)"
    )
    
    # SLIDE 19: Testing & Validation
    add_content_slide(
        "Testing & Validation",
        [
            "Multi-Layer Testing Strategy",
            "",
            "UNIT TESTS (Base - Most Tests):",
            "  • merkle_tree_service_test.dart",
            "  • blockchain_service_test.dart",
            "  • auth_service_test.dart",
            "",
            "INTEGRATION TESTS (Middle):",
            "  • Firestore Security Rules validation",
            "  • Cloud Functions local emulation",
            "  • API endpoint testing (Gemini AI)",
            "",
            "UI TESTS (Widget tests):",
            "  • Authentication screens",
            "  • QR code generation/scanning",
            "  • AI chat interface",
            "",
            "MANUAL TESTING (Top - Critical):",
            "  • End-to-end user journeys",
            "  • Biometric lock on real devices",
            "  • Blockchain transaction on testnet",
            "",
            "TEST COVERAGE: 70% overall"
        ],
        notes="Tested with 5 beta users, 50+ documents, 15+ blockchain transactions"
    )
    
    # SLIDE 20: Challenges & Solutions
    add_content_slide(
        "Challenges & Solutions",
        [
            "CHALLENGE: Biometric lock triggered infinitely",
            "SOLUTION: Static state persistence across rebuilds",
            "",
            "CHALLENGE: Firestore quota exceeded during AI queries",
            "SOLUTION: Intelligent record selection (50 most relevant)",
            "",
            "CHALLENGE: Polygon gas price below minimum error",
            "SOLUTION: Fetch current gas prices, enforce 30 Gwei minimum",
            "",
            "CHALLENGE: OTP expired before user enters code",
            "SOLUTION: 5-minute countdown timer + resend button",
            "",
            "CHALLENGE: PDF parsing failed on some documents",
            "SOLUTION: Cloud Function fallback with pdf-parse library",
            "",
            "LEARNING OUTCOMES:",
            "  • Blockchain has hidden complexities",
            "  • LLM prompts require extensive iteration (50+ revisions)",
            "  • Mobile app lifecycle management is nuanced"
        ],
        notes="Share one challenge in detail with team problem-solving narrative"
    )
    
    # SLIDE 21: Future Enhancements
    add_content_slide(
        "Future Enhancements - Roadmap",
        [
            "PHASE 1 (Months 1-2): Production Readiness",
            "  • Migrate to Polygon Mainnet",
            "  • Implement full HIPAA compliance",
            "  • Add encrypted backups",
            "  • iOS App Store + Google Play deployment",
            "",
            "PHASE 2 (Months 3-4): Feature Expansion",
            "  • Telemedicine integration (video consultations)",
            "  • Lab result auto-import (API integrations)",
            "  • Medication reminders (push notifications)",
            "  • Family account linking (parent-child records)",
            "",
            "PHASE 3 (Months 5-6): Advanced AI",
            "  • Multimodal Gemini (analyze X-rays, MRIs)",
            "  • Predictive health insights",
            "  • Symptom checker chatbot",
            "  • Drug interaction warnings",
            "",
            "LONG-TERM VISION:",
            "  • Government healthcare integration (ABDM)",
            "  • Insurance claim automation",
            "  • Clinical trial recruitment"
        ],
        notes="Emphasize focus on India first (regulatory landscape, market need)"
    )
    
    # SLIDE 22: Compliance & Regulations
    add_content_slide(
        "Compliance & Regulations",
        [
            "✅ IMPLEMENTED:",
            "  • Data encryption (AES-256, TLS)",
            "  • Access control (role-based, audit logged)",
            "  • Patient consent mechanism (QR-based)",
            "  • Right to be Forgotten (account deletion)",
            "  • Minimal data collection",
            "",
            "⚠️ PARTIALLY IMPLEMENTED:",
            "  • HIPAA Compliance (technical controls ✅, legal BAA ⏳)",
            "  • GDPR Compliance (data privacy ✅, EU hosting ⏳)",
            "",
            "❌ FUTURE WORK:",
            "  • India Digital Personal Data Protection Act 2023",
            "  • Clinical validation studies",
            "  • Third-party security audit (penetration testing)",
            "",
            "LEGAL DISCLAIMER:",
            "  • ChainCare is a medical records repository, NOT diagnostic",
            "  • AI assistants are educational, not medical advice",
            "  • All clinical decisions must involve licensed physicians"
        ],
        notes="Academic prototype; production requires legal consultation"
    )
    
    # SLIDE 23: Competitive Analysis
    add_content_slide(
        "Competitive Analysis",
        [
            "ChainCare vs Competitors",
            "",
            "                    ChainCare | Practo | Apollo 247 | Google Health",
            "Digital Records        ✅     |   ✅   |     ✅     |      ✅",
            "Blockchain Audit       ✅     |   ❌   |     ❌     |      ❌",
            "AI Assistant (Dual)    ✅     |   ❌   |     ❌     |      ⚠️",
            "QR-Based Consent       ✅     |   ❌   |     ❌     |      ❌",
            "OCR Upload             ✅     |   ❌   |     ✅     |      ❌",
            "Open Source            ✅     |   ❌   |     ❌     |      ❌",
            "Offline Support        ✅     |   ⚠️   |     ⚠️     |      ❌",
            "",
            "UNIQUE SELLING POINTS:",
            "  1. ONLY blockchain-verified audit trail in India",
            "  2. ONLY dual AI assistants (doctor + patient)",
            "  3. ONLY explicit QR-based consent mechanism",
            "  4. Open-source academic project"
        ],
        notes="Acknowledge Practo/Apollo have telemedicine (we don't), but we have blockchain"
    )
    
    # SLIDE 24: Project Timeline
    add_content_slide(
        "Project Timeline",
        [
            "6-Month Development Journey",
            "",
            "MONTH 1 (July 2025): Planning & Research",
            "  Technology stack, Firebase setup, UI/UX mockups",
            "",
            "MONTH 2 (August 2025): Core Development",
            "  Authentication flows, document upload, database design",
            "",
            "MONTH 3 (September 2025): Advanced Features",
            "  Biometric integration, ML Kit OCR, blockchain, AI",
            "",
            "MONTH 4 (October 2025): Integration & Testing",
            "  Cross-module integration, security rules, beta testing",
            "",
            "MONTH 5 (November 2025): Polish & Optimization",
            "  UI/UX refinements, performance optimization, bug fixes",
            "",
            "MONTH 6 (December 2025): Final Prep",
            "  Presentation creation, live demo, documentation",
            "",
            "MILESTONES:",
            "  ✅ MVP (end of Month 2)",
            "  ✅ Feature-complete (end of Month 4)",
            "  ✅ Production-ready (end of Month 6)"
        ],
        notes="Highlight parallel development (4 members working simultaneously)"
    )
    
    # SLIDE 25: Team Contributions
    add_content_slide(
        "Team Contributions",
        [
            "Balanced Module Ownership",
            "",
            "[Member 1]: Authentication & User Management",
            "  Files: 11 | LOC: 2,400 | Hours: 35",
            "  Features: Email/Google/Phone auth, Biometric lock, Profiles",
            "",
            "[Member 2]: Medical Records & Data Management",
            "  Files: 8 | LOC: 2,390 | Hours: 40",
            "  Features: Document upload, ML Kit OCR, QR system",
            "",
            "[Member 3]: Blockchain Audit Trail",
            "  Files: 6 | LOC: 2,510 | Hours: 45",
            "  Features: Smart contract, Merkle trees, Cloud Functions",
            "",
            "[Member 4]: AI Clinical Assistants",
            "  Files: 10 | LOC: 4,850 | Hours: 45",
            "  Features: RAG architecture, Dual AI, Admin dashboard",
            "",
            "TOTAL: 165 hours | 9,500+ lines of code",
            "",
            "Weekly sync meetings for cross-module testing"
        ],
        notes="Each member presents their module for 5-7 minutes during this section"
    )
    
    # SLIDE 26: Live Demo
    add_content_slide(
        "Live Demo",
        [
            "Demo Flow (5-10 minutes)",
            "",
            "1. PATIENT JOURNEY (3 min):",
            "   • Launch app → Biometric lock",
            "   • Login → Patient dashboard",
            "   • Upload medical document",
            "   • Generate QR Medical ID",
            "   • Chat with AI: 'What is my hemoglobin level?'",
            "",
            "2. DOCTOR JOURNEY (3 min):",
            "   • Login as doctor",
            "   • Scan patient QR → Consent granted",
            "   • View patient records",
            "   • Chat with Clinical Assistant",
            "",
            "3. ADMIN JOURNEY (2 min):",
            "   • View pending doctor applications",
            "   • Check ML confidence score",
            "   • Approve doctor",
            "",
            "4. BLOCKCHAIN VERIFICATION (2 min):",
            "   • View transaction on PolygonScan",
            "   • Verify Merkle root matches"
        ],
        notes="Backup plan: Pre-recorded video if live demo fails"
    )
    
    # SLIDE 27: Technical Achievements
    add_content_slide(
        "Technical Achievements",
        [
            "Notable Technical Accomplishments",
            "",
            "🏆 FULL-STACK DEVELOPMENT",
            "   Flutter mobile app (9,500 LOC)",
            "   Node.js backend (600 LOC)",
            "   Solidity smart contract (210 LOC)",
            "",
            "🏆 BLOCKCHAIN INTEGRATION",
            "   15+ successful transactions on Polygon Amoy",
            "   Merkle tree implementation from scratch",
            "",
            "🏆 AI/ML IMPLEMENTATION",
            "   RAG architecture for context-aware AI",
            "   700+ lines of prompt engineering",
            "",
            "🏆 SECURITY BEST PRACTICES",
            "   6-layer defense-in-depth",
            "   Server-side Firestore Security Rules",
            "",
            "🏆 PERFORMANCE OPTIMIZATION",
            "   28% app size reduction",
            "   84% image compression",
            "",
            "🏆 PRODUCTION-GRADE CODE",
            "   Comprehensive error handling",
            "   Logging and monitoring"
        ],
        notes="Position as production-ready, not just academic proof-of-concept"
    )
    
    # SLIDE 28: Learnings & Takeaways
    add_content_slide(
        "Learnings & Takeaways",
        [
            "Key Learnings from 6-Month Journey",
            "",
            "💡 TECHNICAL INSIGHTS:",
            "  • Blockchain requires deep protocol understanding",
            "  • LLMs need extensive prompt iteration (50+ revisions)",
            "  • Mobile app lifecycle is complex",
            "  • Firebase free tier is generous but requires design foresight",
            "",
            "💡 TEAMWORK:",
            "  • Clear module boundaries enable parallel development",
            "  • Weekly integration meetings prevent conflicts",
            "  • Documentation is critical for team collaboration",
            "",
            "💡 PROBLEM-SOLVING:",
            "  • Always have a backup plan",
            "  • Error messages should be user-friendly",
            "  • Testing on real devices reveals hidden issues",
            "",
            "💡 HEALTHCARE DOMAIN:",
            "  • Medical jargon intimidating for patients",
            "  • Doctors value precision over friendliness",
            "  • Regulatory compliance is complex"
        ],
        notes="Share specific 'aha moment' or breakthrough during development"
    )
    
    # SLIDE 29: Conclusion & Impact
    add_content_slide(
        "Conclusion & Impact",
        [
            "ChainCare: Bridging Healthcare & Technology",
            "",
            "WHAT WE BUILT:",
            "  ✅ Production-grade mobile app (iOS + Android)",
            "  ✅ Blockchain-secured audit trail (Polygon)",
            "  ✅ Dual AI assistants (Doctor + Patient)",
            "  ✅ Comprehensive security (6 layers)",
            "",
            "REAL-WORLD IMPACT:",
            "  For Patients: Centralized records, AI education, data control",
            "  For Doctors: AI insights, verified credentials, transparency",
            "  For Healthcare: Tamper-proof records, reduced fraud",
            "",
            "BY THE NUMBERS:",
            "  • 9,500+ lines of code",
            "  • 6 cutting-edge technologies",
            "  • 4 team members collaborating",
            "  • 165 hours of development",
            "",
            "VISION:",
            "  'Empowering patients and doctors with secure, intelligent,",
            "   and transparent medical record management through",
            "   blockchain and AI.'"
        ],
        notes="End on inspirational note about technology improving healthcare accessibility"
    )
    
    # SLIDE 30: Q&A
    add_content_slide(
        "Questions & Answers",
        [
            "Thank You!",
            "",
            "We're ready for your questions.",
            "",
            "",
            "ANTICIPATED QUESTIONS:",
            "",
            "Q: Why Polygon instead of Ethereum?",
            "A: Cost - Polygon $0.004/tx vs Ethereum $50-200/tx",
            "",
            "Q: How do you handle HIPAA compliance?",
            "A: Technical controls implemented, legal BAA needed for production",
            "",
            "Q: Can blockchain be hacked?",
            "A: Polygon Mainnet has $7B+ locked, extremely secure",
            "",
            "Q: What if Firebase goes down?",
            "A: Offline persistence - app works from cache",
            "",
            "Q: How accurate is the AI?",
            "A: Educational tool, NOT diagnostic. Doctors make all decisions."
        ],
        notes="Stay calm, if unsure say 'Great question for further research' rather than guessing"
    )
    
    # Save presentation
    output_file = "ChainCare_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created successfully: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    create_chaincare_presentation()
